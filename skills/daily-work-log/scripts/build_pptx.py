"""
A4 PPT 생성 (python-pptx). HTML 템플릿과 시각적으로 동일하게 렌더링.
흰 바탕 + 라벨/헤더는 연회색(#f2f2f2) + 회색 셀 보더(#808080) + 업무명/마지막 행 굵은 검정 보더.
기본 PowerPoint 테이블 스타일(파란 헤더 등)은 제거한다.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


STATUS_W = Mm(210)
STATUS_H = Mm(297)
GANTT_W = Mm(297)
GANTT_H = Mm(210)

MARGIN_L = Mm(12)
MARGIN_R = Mm(12)
MARGIN_T = Mm(12)
MARGIN_B = Mm(10)

FONT_FACE = 'Malgun Gothic'

COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LABEL_BG = RGBColor(242, 242, 242)   # #f2f2f2
COLOR_BORDER_SOFT = RGBColor(128, 128, 128)  # #808080
COLOR_FOOTER = RGBColor(85, 85, 85)

# 보더 두께 (EMU, 1pt = 12700)
WIDTH_THIN = 6350    # ~0.5pt (HTML 0.5px border)
WIDTH_THICK = 15240  # ~1.2pt (HTML 1.2px accent border)


def _a4_prs(orientation='portrait'):
    prs = Presentation()
    if orientation == 'landscape':
        prs.slide_width = GANTT_W
        prs.slide_height = GANTT_H
    else:
        prs.slide_width = STATUS_W
        prs.slide_height = STATUS_H
    return prs


def _set_run(run, text, *, size=Pt(9), bold=False, color=COLOR_BLACK, face=FONT_FACE):
    run.text = text
    run.font.name = face
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _clear_paragraph(p):
    for r in list(p.runs):
        r._r.getparent().remove(r._r)


def _clear_table_style(tbl):
    """기본 PPT 테이블 스타일(파란 헤더·banding) 제거 → HTML과 동일한 흰 바탕으로."""
    tbl_elem = tbl._tbl
    tblPr = tbl_elem.find(qn('a:tblPr'))
    if tblPr is None:
        return
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')
    tblPr.set('firstCol', '0')
    tblPr.set('bandCol', '0')
    for child in list(tblPr):
        if child.tag == qn('a:tableStyleId'):
            tblPr.remove(child)


def _set_cell_text(cell, text_or_lines, *, size=Pt(9), bold=False, color=COLOR_BLACK,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bullet=False, bg=COLOR_WHITE):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Mm(2.2)
    tf.margin_right = Mm(1.5)
    tf.margin_top = Mm(1.0)
    tf.margin_bottom = Mm(1.0)
    cell.vertical_anchor = anchor
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg

    if isinstance(text_or_lines, str):
        lines = [text_or_lines]
    else:
        lines = list(text_or_lines) if text_or_lines else [""]

    for _ in range(len(tf.paragraphs) - 1):
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)

    first = tf.paragraphs[0]
    _clear_paragraph(first)
    first.alignment = align

    def _add_text(p, txt):
        run = p.add_run()
        prefix = '· ' if bullet else ''
        _set_run(run, prefix + txt, size=size, bold=bold, color=color)

    _add_text(first, lines[0])
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        _add_text(p, line)


def _set_border(cell, side, color=COLOR_BORDER_SOFT, width_emu=WIDTH_THIN):
    """side: 'T','B','L','R'"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tag = f'a:ln{side}'
    ln = tcPr.find(qn(tag))
    if ln is None:
        ln = etree.SubElement(tcPr, qn(tag))
    ln.set('w', str(width_emu))
    ln.set('cap', 'flat')
    ln.set('cmpd', 'sng')
    ln.set('algn', 'ctr')
    for child in list(ln):
        ln.remove(child)
    sf = etree.SubElement(ln, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
    pd = etree.SubElement(ln, qn('a:prstDash'))
    pd.set('val', 'solid')
    etree.SubElement(ln, qn('a:round'))


def _all_borders(cell, color=COLOR_BORDER_SOFT, width_emu=WIDTH_THIN):
    for side in ('T', 'B', 'L', 'R'):
        _set_border(cell, side, color, width_emu)


def _row_height(text_lines, *, per_line_mm=4.0, min_mm=7.0):
    n = len(text_lines) if isinstance(text_lines, (list, tuple)) else 1
    return Mm(max(min_mm, n * per_line_mm + 2.0))


# ---------- Status PPT ----------

def build_status_pptx(out_path: Path, categories: list, year_short: str, month: int, week_num: int, author: str = ''):
    prs = _a4_prs('portrait')
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    content_w = STATUS_W - MARGIN_L - MARGIN_R
    col_w = content_w

    # Title
    tb = slide.shapes.add_textbox(MARGIN_L, MARGIN_T, Mm(150), Mm(9))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    _clear_paragraph(tf.paragraphs[0])
    _set_run(tf.paragraphs[0].add_run(), '주간 업무 현황 _ 기획', size=Pt(16), bold=True)

    # Week label
    tb2 = slide.shapes.add_textbox(STATUS_W - MARGIN_R - Mm(50), MARGIN_T + Mm(1.5), Mm(50), Mm(7))
    tf2 = tb2.text_frame
    tf2.margin_left = tf2.margin_right = Emu(0)
    tf2.margin_top = tf2.margin_bottom = Emu(0)
    _clear_paragraph(tf2.paragraphs[0])
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _set_run(tf2.paragraphs[0].add_run(),
             f"’{year_short}년 {month}월 {week_num}주차", size=Pt(10), bold=True)

    body_top = MARGIN_T + Mm(10)
    body_h = STATUS_H - body_top - MARGIN_B

    # 카테고리별 테이블 구조
    tables_info = []
    for cat in categories:
        rows = [
            ('task', [cat['name']], _row_height(cat['name'], per_line_mm=4, min_mm=6.5)),
            ('overview', [cat['overview']], _row_height(cat['overview'], per_line_mm=3.8, min_mm=6.0)),
            ('done', cat['done'], _row_height(cat['done'], per_line_mm=4.0, min_mm=6.0)),
        ]
        if cat.get('plan'):
            rows.append(('plan', cat['plan'], _row_height(cat['plan'], per_line_mm=4.0, min_mm=6.0)))
        tables_info.append(rows)

    gap = Mm(2.0)
    total_rows_h = sum(sum(r[2] for r in tbl) for tbl in tables_info)
    total_gaps = gap * max(0, len(tables_info) - 1)
    scale = body_h / (total_rows_h + total_gaps) if (total_rows_h + total_gaps) > body_h else 1.0

    current_top = body_top
    for rows in tables_info:
        n_rows = len(rows)
        heights = [int(r[2] * scale) for r in rows]
        table_h = sum(heights)
        tbl = slide.shapes.add_table(n_rows, 2, MARGIN_L, current_top, col_w, table_h).table
        _clear_table_style(tbl)

        label_col_w = Mm(17)
        tbl.columns[0].width = label_col_w
        tbl.columns[1].width = col_w - label_col_w
        for i, h in enumerate(heights):
            tbl.rows[i].height = h

        for r_idx, (kind, content, _h) in enumerate(rows):
            lc = tbl.cell(r_idx, 0)
            vc = tbl.cell(r_idx, 1)
            is_first = (r_idx == 0)
            is_last = (r_idx == n_rows - 1)

            label_text = {'task': '업무명', 'overview': '개요', 'done': '실적', 'plan': '계획'}[kind]
            _set_cell_text(lc, label_text, size=Pt(8.5), bold=True,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_LABEL_BG)

            if kind == 'task':
                _set_cell_text(vc, content, size=Pt(10), bold=True,
                               anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_WHITE)
            elif kind == 'overview':
                _set_cell_text(vc, content, size=Pt(8.5),
                               anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_WHITE)
            else:
                _set_cell_text(vc, content, size=Pt(8.8), bullet=True,
                               anchor=MSO_ANCHOR.TOP, bg=COLOR_WHITE)

            # 기본 얇은 회색 보더
            _all_borders(lc, COLOR_BORDER_SOFT, WIDTH_THIN)
            _all_borders(vc, COLOR_BORDER_SOFT, WIDTH_THIN)
            # task-row 상단 + last-row 하단: 굵은 검정
            if is_first:
                _set_border(lc, 'T', COLOR_BLACK, WIDTH_THICK)
                _set_border(vc, 'T', COLOR_BLACK, WIDTH_THICK)
            if is_last:
                _set_border(lc, 'B', COLOR_BLACK, WIDTH_THICK)
                _set_border(vc, 'B', COLOR_BLACK, WIDTH_THICK)

        current_top = current_top + table_h + gap

    # Footer (author)
    if author:
        foot = slide.shapes.add_textbox(STATUS_W - MARGIN_R - Mm(30), STATUS_H - Mm(7), Mm(30), Mm(4))
        tf = foot.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        _clear_paragraph(tf.paragraphs[0])
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _set_run(tf.paragraphs[0].add_run(), author, size=Pt(8.5), color=COLOR_FOOTER)

    prs.save(str(out_path))
    return out_path


# ---------- Gantt PPT ----------

def build_gantt_pptx(out_path: Path, sheets: list, author: str = ''):
    prs = _a4_prs('landscape')
    GANTT_ML = Mm(10)
    GANTT_MR = Mm(10)
    GANTT_MT = Mm(9)
    GANTT_MB = Mm(8)

    for sheet in sheets:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        content_w = GANTT_W - GANTT_ML - GANTT_MR

        # Title
        tb = slide.shapes.add_textbox(GANTT_ML, GANTT_MT, Mm(200), Mm(9))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        _clear_paragraph(tf.paragraphs[0])
        _set_run(tf.paragraphs[0].add_run(), sheet['title'], size=Pt(16), bold=True)

        # Author
        if author:
            tb2 = slide.shapes.add_textbox(GANTT_W - GANTT_MR - Mm(40), GANTT_MT + Mm(1.5), Mm(40), Mm(7))
            tf2 = tb2.text_frame
            tf2.margin_left = tf2.margin_right = Emu(0)
            tf2.margin_top = tf2.margin_bottom = Emu(0)
            _clear_paragraph(tf2.paragraphs[0])
            tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
            _set_run(tf2.paragraphs[0].add_run(), author, size=Pt(10), bold=True)

        table_top = GANTT_MT + Mm(11)
        table_h = GANTT_H - table_top - GANTT_MB
        n_rows = 1 + len(sheet['rows'])
        n_cols = 7

        tbl = slide.shapes.add_table(n_rows, n_cols, GANTT_ML, table_top, content_w, table_h).table
        _clear_table_style(tbl)

        # 컬럼 너비 (HTML과 동일 비율: 업무 10%, 일정 6%, 요일 16.8%×5)
        task_w = int(content_w * 0.10)
        due_w = int(content_w * 0.06)
        day_w = int((content_w - task_w - due_w) / 5)
        tbl.columns[0].width = task_w
        tbl.columns[1].width = due_w
        for i in range(2, 7):
            tbl.columns[i].width = day_w

        header_h = Mm(10)
        data_h = (table_h - header_h) // len(sheet['rows'])
        tbl.rows[0].height = header_h
        for i in range(1, n_rows):
            tbl.rows[i].height = data_h

        # Header row: 연회색 배경 + 검정 글자 + 검정 얇은 보더 (HTML과 동일)
        headers = [
            '업무', '일정',
            f"월\n{sheet['dates']['mon']}",
            f"화\n{sheet['dates']['tue']}",
            f"수\n{sheet['dates']['wed']}",
            f"목\n{sheet['dates']['thu']}",
            f"금\n{sheet['dates']['fri']}",
        ]
        for c_idx, h in enumerate(headers):
            cell = tbl.cell(0, c_idx)
            lines = h.split('\n')
            _set_cell_text(cell, lines, size=Pt(9.5), bold=True, color=COLOR_BLACK,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_LABEL_BG)
            _all_borders(cell, COLOR_BLACK, WIDTH_THIN)

        # Data rows: 흰 바탕 (banding 없음). 업무/일정 컬럼만 연회색.
        day_keys = ['mon', 'tue', 'wed', 'thu', 'fri']
        for r_idx, row in enumerate(sheet['rows'], start=1):
            # 업무
            c = tbl.cell(r_idx, 0)
            _set_cell_text(c, row['task'], size=Pt(9), bold=True,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_LABEL_BG)
            _all_borders(c, COLOR_BORDER_SOFT, WIDTH_THIN)
            # 일정
            c = tbl.cell(r_idx, 1)
            _set_cell_text(c, row.get('due', ''), size=Pt(9), bold=True, color=COLOR_BLACK,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bg=COLOR_LABEL_BG)
            _all_borders(c, COLOR_BORDER_SOFT, WIDTH_THIN)
            # 요일 셀
            for c_idx, key in enumerate(day_keys, start=2):
                items = row.get(key, [])
                c = tbl.cell(r_idx, c_idx)
                _set_cell_text(c, items if items else "", size=Pt(8.5), bullet=bool(items),
                               anchor=MSO_ANCHOR.TOP, bg=COLOR_WHITE)
                _all_borders(c, COLOR_BORDER_SOFT, WIDTH_THIN)

    prs.save(str(out_path))
    return out_path
